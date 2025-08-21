import streamlit as st
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import plotly.express as px

st.set_page_config(page_title="Export Report", page_icon="🖨️", layout="wide")
st.title("🖨️ Export PowerPoint Report")

df = st.session_state.get("df")
if df is None or df.empty:
    st.warning("No data available. Go back to the Home page to upload and map your MIS file.")
    st.stop()

metric = "revenue" if "revenue" in df.columns else ("gmv" if "gmv" in df.columns else None)

def kpi_box(title, value):
    st.metric(title, value)

# Preview KPIs
st.subheader("KPIs Preview")
col = st.columns(4)
vals = {}
with col[0]:
    vals["Revenue"] = float(df.get(metric, pd.Series(dtype=float)).sum()) if metric else 0.0
    kpi_box("Revenue", f"{vals['Revenue']:,.0f}")
with col[1]:
    vals["Orders"] = float(df.get("orders", pd.Series(dtype=float)).sum())
    kpi_box("Orders", f"{vals['Orders']:,.0f}")
with col[2]:
    vals["Units"] = float(df.get("units", pd.Series(dtype=float)).sum())
    kpi_box("Units", f"{vals['Units']:,.0f}")
with col[3]:
    if "cost" in df.columns and metric:
        gm = (df[metric] - df["cost"]).sum()
        vals["GM %"] = round(100*(1 - df["cost"].sum()/ (df[metric].sum() or 1)), 2)
        kpi_box("GM %", f"{vals['GM %']}%")
    else:
        vals["GM %"] = None
        kpi_box("GM %", "—")

# Build a couple of figures to embed
figs = []
if "date" in df.columns and metric:
    ts = df.groupby("date", as_index=False)[metric].sum().sort_values("date")
    figs.append(px.line(ts, x="date", y=metric, markers=True, title="Revenue over time"))
if "region" in df.columns and metric:
    reg = df.groupby("region", as_index=False)[metric].sum().sort_values(metric, ascending=False).head(10)
    figs.append(px.bar(reg, x="region", y=metric, title="Top Regions by Revenue"))
if "product" in df.columns and metric:
    prod = df.groupby("product", as_index=False)[metric].sum().sort_values(metric, ascending=False).head(10)
    figs.append(px.bar(prod, x="product", y=metric, title="Top Products by Revenue"))

def fig_to_image_bytes(fig):
    # Using plotly static export via to_image requires kaleido; instead, use HTML image export workaround.
    # To keep deps light, we'll export tables and KPIs; for figures we add placeholders text.
    return None

def build_ppt(df: pd.DataFrame) -> BytesIO:
    prs = Presentation()
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "MIS Dashboard Report"
    slide.placeholders[1].text = "Auto-generated from Streamlit • Powered by Python"

    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Key Metrics"
    left, top = Inches(0.5), Inches(1.5)
    tx_box = slide.shapes.add_textbox(left, top, Inches(9), Inches(4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    lines = [
        f"Revenue: {vals['Revenue']:,.0f}",
        f"Orders: {vals['Orders']:,.0f}",
        f"Units: {vals['Units']:,.0f}",
        f"GM %: {vals['GM %']}%" if vals['GM %'] is not None else "GM %: —",
    ]
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(24)

    # Top tables slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top Breakdowns"
    metric_col = metric if metric else "revenue"

    y = Inches(1.5)
    def add_table(title, df_top, y):
        title_shape = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.4))
        title_shape.text_frame.text = title
        title_shape.text_frame.paragraphs[0].font.bold = True
        y += Inches(0.3)
        rows, cols = min(6, len(df_top)+1), len(df_top.columns)
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), y, Inches(9), Inches(1.5))
        table = table_shape.table
        for j, col in enumerate(df_top.columns):
            table.cell(0, j).text = str(col)
        for i in range(min(5, len(df_top))):
            for j, col in enumerate(df_top.columns):
                table.cell(i+1, j).text = str(df_top.iloc[i, j])
        return y + Inches(1.9)

    if "region" in df.columns and metric:
        reg = df.groupby("region", as_index=False)[metric_col].sum().sort_values(metric_col, ascending=False).head(5)
        y = add_table("Top Regions", reg, y)
    if "product" in df.columns and metric:
        prod = df.groupby("product", as_index=False)[metric_col].sum().sort_values(metric_col, ascending=False).head(5)
        y = add_table("Top Products", prod, y)
    if "channel" in df.columns and metric:
        ch = df.groupby("channel", as_index=False)[metric_col].sum().sort_values(metric_col, ascending=False).head(5)
        y = add_table("Top Channels", ch, y)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

if st.button("Generate PowerPoint"):
    bio = build_ppt(df)
    st.download_button("Download Report.pptx", data=bio, file_name="MIS_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
else:
    st.info("Click **Generate PowerPoint** to build a deck based on the current filtered view.")